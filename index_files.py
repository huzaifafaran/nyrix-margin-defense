import os
import pandas as pd
from pptx import Presentation
import glob

def extract_pptx_content(filepath):
    """Extracts text content from a PPTX file."""
    try:
        prs = Presentation(filepath)
        content = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            if slide.shapes.title:
                slide_content.append(f"Title: {slide.shapes.title.text}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
            content.append(f"Slide {i+1}:\n" + "\n".join(slide_content))
        return "\n---\n".join(content)
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def extract_xlsx_content(filepath):
    """Extracts a summary of an XLSX file."""
    try:
        xl = pd.ExcelFile(filepath)
        summary = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            summary.append(f"Sheet: {sheet_name}")
            summary.append("Columns: " + ", ".join(df.columns.astype(str)))
            summary.append(f"Rows: {len(df)}")
            summary.append("First 5 rows preview:")
            summary.append(df.head().to_markdown(index=False))
            summary.append("-" * 20)
        return "\n".join(summary)
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def main():
    base_dir = r"c:\Users\Gigabyte\Desktop\lucky cement testing"
    output_file = "file_index.md"
    
    with open(output_file, "w", encoding="utf-8") as f:
        f.write("# File Index\n\n")
        
        # PPTX Files
        for filepath in glob.glob(os.path.join(base_dir, "*.pptx")):
            filename = os.path.basename(filepath)
            print(f"Processing {filename}...")
            f.write(f"## {filename}\n\n")
            content = extract_pptx_content(filepath)
            f.write(content + "\n\n")
            
        # XLSX Files
        for filepath in glob.glob(os.path.join(base_dir, "*.xlsx")):
            filename = os.path.basename(filepath)
            print(f"Processing {filename}...")
            f.write(f"## {filename}\n\n")
            content = extract_xlsx_content(filepath)
            f.write(content + "\n\n")

    print(f"Indexing complete. Saved to {output_file}")

if __name__ == "__main__":
    main()
